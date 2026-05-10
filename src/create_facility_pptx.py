#!/usr/bin/env python3
"""
施設レポート PowerPoint 自動生成スクリプト

Usage:
    python src/create_facility_pptx.py \\
        --facility "施設名" \\
        --period "2026年1月〜3月" \\
        --occupancy data/sample/occupancy.csv \\
        --changes data/sample/changes.csv \\
        --reservations data/sample/reservations.csv \\
        --output report.pptx
"""

import argparse
import csv
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm

# Use IPA Gothic for Japanese characters in charts if available
_JP_FONT_PATH = "/usr/share/fonts/opentype/ipafont-gothic/ipag.ttf"
if Path(_JP_FONT_PATH).exists():
    fm.fontManager.addfont(_JP_FONT_PATH)
    plt.rcParams["font.family"] = "IPAGothic"

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

COLOR_ACCENT = RGBColor(0x1F, 0x49, 0x7D)
COLOR_LIGHT = RGBColor(0xD9, 0xE2, 0xF3)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_HIGHLIGHT = RGBColor(0xFF, 0xFF, 0xCC)
COLOR_DARK = RGBColor(0x40, 0x40, 0x40)
COLOR_GRAY = RGBColor(0x80, 0x80, 0x80)


# ---------------------------------------------------------------------------
# CSV loaders
# ---------------------------------------------------------------------------

def load_occupancy(path):
    data = []
    with open(path, newline="", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            data.append({"month": row["month"], "rate": float(row["rate"])})
    return data


def load_changes(path):
    data = []
    with open(path, newline="", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            data.append({"date": row["date"], "count": int(row["count"])})
    return data


def load_reservations(path):
    data = []
    with open(path, newline="", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            data.append({
                "month": row["month"],
                "channel": row["channel"],
                "reservations": int(row["reservations"]),
                "rooms": int(row["rooms"]),
            })
    return data


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _add_header_bar(slide, title_text):
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.08), Inches(12), Inches(0.5))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE


def _add_textbox(slide, text, left, top, width, height, size=12, bold=False, color=None):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLOR_DARK


def _set_cell(cell, text, size=10, bold=False,
              fg=None, bg=None, align=PP_ALIGN.CENTER):
    cell.text = text
    tf = cell.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if fg:
        run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_cover(prs, facility, period):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    _add_textbox(slide, "施設レポート",
                 Inches(0.5), Inches(0.35), Inches(12), Inches(0.9),
                 size=30, bold=True, color=COLOR_WHITE)

    _add_textbox(slide, facility,
                 Inches(0.5), Inches(2.0), Inches(12), Inches(1.2),
                 size=40, bold=True, color=COLOR_ACCENT)

    _add_textbox(slide, period,
                 Inches(0.5), Inches(3.5), Inches(12), Inches(0.8),
                 size=24, color=COLOR_DARK)

    _add_textbox(slide, "手間いらず自動　本格運用：2026年1月〜",
                 Inches(0.5), Inches(4.6), Inches(12), Inches(0.6),
                 size=15, color=COLOR_GRAY)


def slide_occupancy_chart(prs, occupancy_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "月別稼働率")

    labels = []
    for d in occupancy_data:
        parts = d["month"].split("-")
        labels.append(f"{int(parts[1])}月" if len(parts) == 2 else d["month"])
    rates = [d["rate"] for d in occupancy_data]

    fig, ax = plt.subplots(figsize=(11, 5.5))
    bars = ax.bar(labels, rates, color="#1F497D", width=0.55)
    ax.set_ylim(0, 100)
    ax.set_ylabel("稼働率 (%)", fontsize=11)
    ax.yaxis.grid(True, linestyle="--", alpha=0.6)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    for bar_, rate in zip(bars, rates):
        ax.text(
            bar_.get_x() + bar_.get_width() / 2,
            bar_.get_height() + 1.0,
            f"{rate:.1f}%",
            ha="center", va="bottom", fontsize=10,
        )

    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)

    slide.shapes.add_picture(buf, Inches(0.5), Inches(0.75), Inches(12.3), Inches(6.4))


def _build_daily_table(slide, month_label, day_rows, top):
    days = sorted(day_rows, key=lambda r: r["date"])
    n = len(days)
    cols = n + 1  # label + day columns

    tbl = slide.shapes.add_table(
        2, cols,
        Inches(0.4), top,
        Inches(12.5), Inches(0.58),
    ).table

    tbl.rows[0].height = Pt(20)
    tbl.rows[1].height = Pt(20)

    _set_cell(tbl.cell(0, 0), month_label,
              size=8, bold=True, fg=COLOR_WHITE, bg=COLOR_ACCENT)
    _set_cell(tbl.cell(1, 0), "変動",
              size=7, bold=True, bg=COLOR_LIGHT)

    for i, d in enumerate(days, start=1):
        day_num = d["date"].split("-")[-1].lstrip("0") or "0"
        count = d["count"]

        _set_cell(tbl.cell(0, i), day_num,
                  size=7, fg=COLOR_WHITE, bg=COLOR_ACCENT)
        _set_cell(tbl.cell(1, i), str(count),
                  size=7,
                  bg=COLOR_HIGHLIGHT if count > 0 else None)


def slide_changes(prs, changes_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "変動回数（2月・3月）")

    feb = [d for d in changes_data if "-02-" in d["date"]]
    mar = [d for d in changes_data if "-03-" in d["date"]]

    _build_daily_table(slide, "2月", feb, Inches(0.85))
    _build_daily_table(slide, "3月", mar, Inches(1.75))


def slide_reservations(prs, reservations_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "予約件数・室数")

    rows = reservations_data
    n = len(rows)
    headers = ["月", "チャネル", "予約件数", "予約室数"]

    tbl = slide.shapes.add_table(
        n + 1, 4,
        Inches(1.5), Inches(0.75),
        Inches(10.3), Inches(0.38 * (n + 1)),
    ).table

    col_widths = [Inches(1.4), Inches(3.8), Inches(2.5), Inches(2.6)]
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = w

    for ci, h in enumerate(headers):
        _set_cell(tbl.cell(0, ci), h,
                  size=11, bold=True,
                  fg=COLOR_WHITE, bg=COLOR_ACCENT)

    aligns = [PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER]
    for ri, row in enumerate(rows, start=1):
        vals = [row["month"], row["channel"],
                str(row["reservations"]), str(row["rooms"])]
        bg = COLOR_LIGHT if ri % 2 == 0 else None
        for ci, (val, align) in enumerate(zip(vals, aligns)):
            _set_cell(tbl.cell(ri, ci), val, size=10, bg=bg, align=align)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="施設レポート PPTX 自動生成")
    parser.add_argument("--facility", required=True, help="施設名")
    parser.add_argument("--period", required=True, help="対象期間（例: 2026年1月〜3月）")
    parser.add_argument("--occupancy", required=True, help="稼働率 CSV パス")
    parser.add_argument("--changes", required=True, help="変動回数 CSV パス")
    parser.add_argument("--reservations", required=True, help="予約件数 CSV パス")
    parser.add_argument("--output", default="report.pptx", help="出力ファイル名")
    args = parser.parse_args()

    occupancy_data = load_occupancy(args.occupancy)
    changes_data = load_changes(args.changes)
    reservations_data = load_reservations(args.reservations)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs, args.facility, args.period)
    slide_occupancy_chart(prs, occupancy_data)
    slide_changes(prs, changes_data)
    slide_reservations(prs, reservations_data)

    prs.save(args.output)
    print(f"保存しました: {args.output}")


if __name__ == "__main__":
    main()
