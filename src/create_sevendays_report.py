#!/usr/bin/env python3
"""
セブンデイズホテル 自動導入経過報告 PPTX 生成
参考：南福岡グリーンホテル向け報告書フォーマット

Usage:
    python src/create_sevendays_report.py \
        --revpar1 <RevPAR_202412202511.csv> \
        --revpar2 <RevPAR_202505202604.csv> \
        --changes <変動回数.csv> \
        --output  sevendays_report.pptx
"""

import argparse
import csv
import io
from collections import defaultdict
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Japanese font
_JP = "/usr/share/fonts/opentype/ipafont-gothic/ipag.ttf"
if Path(_JP).exists():
    fm.fontManager.addfont(_JP)
    plt.rcParams["font.family"] = "IPAGothic"

SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)
FACILITY = "セブンデイズホテル"
ROOM     = "デイズ S"

# Colors matching reference PDF
C_NAVY  = RGBColor(0x1F, 0x38, 0x64)
C_LBLUE = RGBColor(0xBD, 0xD7, 0xEE)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK  = RGBColor(0x25, 0x25, 0x25)
C_GRAY  = RGBColor(0x70, 0x70, 0x70)
C_NOTE  = RGBColor(0xD6, 0xE4, 0xF0)
C_RED   = RGBColor(0xC0, 0x00, 0x00)

NAVY_HEX  = "#1F3864"
BLUE_HEX  = "#4472C4"
RED_HEX   = "#C00000"
LBLUE_HEX = "#BDD7EE"


# ---------------------------------------------------------------------------
# CSV loaders
# ---------------------------------------------------------------------------

def load_revpar(path):
    rows = []
    with open(path, newline="", encoding="shift_jis", errors="replace") as f:
        reader = csv.reader(f)
        next(reader)
        for row in reader:
            if not row or not row[0].strip():
                continue
            rows.append({
                "month":   row[0][:7],
                "rooms":   int(row[1]),
                "revenue": int(row[2]),
                "adr":     float(row[3]),
                "occ":     float(row[4]),
                "revpar":  float(row[5]),
            })
    return rows


def load_changes(path):
    """Returns {year: {day_index: total_count}}  (day_index 0 = Feb 1)."""
    totals = defaultdict(lambda: defaultdict(int))
    with open(path, newline="", encoding="shift_jis", errors="replace") as f:
        reader = csv.reader(f)
        next(reader)
        for row in reader:
            if not row or not row[0].strip():
                continue
            year = row[1].strip()[:4]
            for i, v in enumerate(row[2:]):
                try:
                    totals[year][i] += int(v.strip())
                except ValueError:
                    pass
    return totals


def daily_for_month(changes_data, year, month):
    """Return [(day_num, count), ...] for the given year/month."""
    yd = changes_data.get(str(year), {})
    if month == 2:
        return [(i + 1,        yd.get(i, 0)) for i in range(0, 28)]
    if month == 3:
        return [(i - 27,       yd.get(i, 0)) for i in range(28, 59)]
    if month == 4:
        return [(i - 58,       yd.get(i, 0)) for i in range(59, 89)]
    return []


# ---------------------------------------------------------------------------
# python-pptx helpers
# ---------------------------------------------------------------------------

def _header_bar(slide, text):
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_NAVY
    bar.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(0.35), Inches(0.1), Inches(12.6), Inches(0.45))
    p   = txb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(20)
    run.font.bold  = True
    run.font.color.rgb = C_WHITE


def _txb(slide, text, left, top, width, height,
         size=12, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color or C_DARK


def _note_box(slide, text, left, top, width, height):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = C_NOTE
    box.line.fill.background()
    txb = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.08),
        width - Inches(0.3), height - Inches(0.12))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = "✅  " + text
    run.font.size  = Pt(12)
    run.font.color.rgb = C_DARK


def _set_cell(cell, text, size=9, bold=False,
              fg=None, bg=None, align=PP_ALIGN.CENTER):
    cell.text = text
    tf = cell.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if fg:
        run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def _fig_buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return buf


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Light gray background
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    bg.line.fill.background()

    # Title band
    band = slide.shapes.add_shape(1, Inches(3.8), Inches(3.0), Inches(9.2), Inches(1.7))
    band.fill.solid()
    band.fill.fore_color.rgb = C_LBLUE
    band.line.fill.background()

    _txb(slide, "自動導入経過報告",
         Inches(3.8), Inches(3.05), Inches(9.2), Inches(0.8),
         size=28, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    _txb(slide, "— 2026年1月〜2026年4月 —",
         Inches(3.8), Inches(3.82), Inches(9.2), Inches(0.6),
         size=17, color=C_DARK, align=PP_ALIGN.CENTER)

    _txb(slide, FACILITY,
         Inches(7.5), Inches(5.1), Inches(5.5), Inches(0.85),
         size=26, bold=True, color=C_NAVY, align=PP_ALIGN.RIGHT)
    _txb(slide, "手間いらず株式会社",
         Inches(7.5), Inches(5.85), Inches(5.5), Inches(0.6),
         size=18, color=C_NAVY, align=PP_ALIGN.RIGHT)


def slide_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(slide, "導入内容")

    items = [
        ("1", "料金在庫管理販売期間",
         "12か月先（365日先）まで、\n24時間体制で「料金」「在庫」を自動監視"),
        ("2", "TP（ターゲットプライス）\n設定ルール",
         "設定ルール\n全部屋 × ランク × 残室数"),
        ("3", "販売チャネル 自動調整",
         "稼働率に応じて海外サイトの売止を\n調整し、収益性を向上"),
    ]
    starts = [Inches(0.4), Inches(4.6), Inches(8.8)]

    for (num, title, body), x in zip(items, starts):
        box = slide.shapes.add_shape(1, x, Inches(1.1), Inches(4.0), Inches(3.6))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFB)
        box.line.color.rgb = C_LBLUE

        nb = slide.shapes.add_shape(1, x + Inches(0.15), Inches(1.25),
                                    Inches(0.42), Inches(0.42))
        nb.fill.solid()
        nb.fill.fore_color.rgb = C_LBLUE
        nb.line.fill.background()
        _txb(slide, num, x + Inches(0.15), Inches(1.25), Inches(0.42), Inches(0.42),
             size=13, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

        _txb(slide, title, x + Inches(0.2), Inches(1.75), Inches(3.6), Inches(0.9),
             size=13, bold=True, color=C_DARK, wrap=True)
        _txb(slide, body, x + Inches(0.2), Inches(2.7), Inches(3.6), Inches(1.8),
             size=11, color=C_DARK, wrap=True)

    _note_box(slide,
              "自動化により日次の料金確認作業を削減しつつ、手動調整との併用で柔軟なレベニュー運用を実現。",
              Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.65))


def _make_changes_fig(daily, month):
    n      = len(daily)
    days   = [d for d, c in daily]
    counts = [c for d, c in daily]
    cumul  = list(np.cumsum(counts))
    half   = 15

    fig = plt.figure(figsize=(12.8, 4.6))
    gs  = fig.add_gridspec(3, 1, height_ratios=[1, 1, 4.2],
                            hspace=0.05, top=0.97, bottom=0.1,
                            left=0.03, right=0.98)

    def _draw_half_table(ax_idx, slice_):
        ax = fig.add_subplot(gs[ax_idx])
        ax.axis("off")
        row_dates  = [f"{month:02d}/{d:02d}" for d, c in slice_]
        row_counts = [str(c)                  for d, c in slice_]
        tbl = ax.table(
            cellText=[row_dates, row_counts],
            rowLabels=["日付", "変動回数"],
            loc="center", cellLoc="center"
        )
        tbl.auto_set_font_size(False)
        tbl.set_fontsize(7.5)
        tbl.scale(1, 1.55)
        for (ri, ci), cell in tbl.get_celld().items():
            cell.set_linewidth(0.4)
            if ci == -1:           # row label
                cell.set_facecolor(NAVY_HEX)
                cell.get_text().set_color("white")
                cell.get_text().set_fontweight("bold")
            elif ri == 0:          # date row
                cell.set_facecolor(LBLUE_HEX)
                cell.get_text().set_fontsize(7)
            else:                  # count row
                cell.set_facecolor("white")

    _draw_half_table(0, daily[:half])
    _draw_half_table(1, daily[half:])

    # Bar + cumulative chart
    ax = fig.add_subplot(gs[2])
    ax.bar(range(1, n + 1), counts,
           color=BLUE_HEX, width=0.7, alpha=0.9, zorder=3, label="変動回数")
    ax.yaxis.grid(True, linestyle="--", alpha=0.4, zorder=0)
    ax.set_axisbelow(True)
    ax.set_xlim(0.5, n + 0.5)
    ax.set_xticks(range(1, n + 1))
    ax.set_xticklabels([f"{month}/{d}" for d in days], fontsize=6.5)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    ax2 = ax.twinx()
    ax2.plot(range(1, n + 1), cumul,
             color=RED_HEX, marker="D", markersize=4,
             linewidth=1.5, label="累計", zorder=4)
    # Annotate select points
    show_idx = set([0, n - 1] + list(range(4, n, 5)))
    for i in show_idx:
        ax2.annotate(str(cumul[i]), (i + 1, cumul[i]),
                     fontsize=7, color=RED_HEX,
                     ha="center", va="bottom",
                     xytext=(0, 4), textcoords="offset points")
    ax2.spines["top"].set_visible(False)
    ax2.tick_params(axis="y", colors=RED_HEX, labelsize=8)

    h1, l1 = ax.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()
    ax.legend(h1 + h2, l1 + l2, loc="upper right", fontsize=8, framealpha=0.8)

    return fig


def slide_changes_month(prs, changes_data, year, month):
    month_jp = {3: "3月", 4: "4月"}[month]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(slide, f"変動履歴対象年月：{year}年{month_jp}")

    daily = daily_for_month(changes_data, year, month)
    total = sum(c for _, c in daily)

    # Big number (top-left)
    _txb(slide, f"{total:,}回",
         Inches(0.4), Inches(0.65), Inches(2.5), Inches(0.85),
         size=42, bold=True, color=C_NAVY)
    _txb(slide, f"※集計期間：2026年1月〜4月\n※対象部屋：{ROOM}（全OTAチャネル合計）",
         Inches(2.9), Inches(0.70), Inches(5.0), Inches(0.75),
         size=9, color=C_GRAY, wrap=True)

    fig = _make_changes_fig(daily, month)
    slide.shapes.add_picture(
        _fig_buf(fig), Inches(0.2), Inches(1.5), Inches(12.9), Inches(5.1))

    _note_box(slide,
              f"TP設定により、{year}年{month_jp}分で {total:,}回 の価格変動を自動実行し、"
              f"機会損失の防止と収益最大化に寄与。",
              Inches(0.3), Inches(6.78), Inches(12.7), Inches(0.58))


def slide_section_break(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1F, 0x6B, 0xBF)
    bg.line.fill.background()
    _txb(slide, title,
         Inches(2), Inches(2.8), Inches(9.33), Inches(1.5),
         size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # TEMAIRAZU logo text
    _txb(slide, "⬤⬤⬤  TEMAIRAZU",
         Inches(4.5), Inches(4.5), Inches(4.3), Inches(0.6),
         size=14, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_monthly_comparison(prs, revpar_all):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(slide, "月次実績比較（チェックイン日ベース　2025年 vs 2026年）")

    by_m = {r["month"]: r for r in revpar_all}
    months   = ["2月", "3月", "4月"]
    keys_25  = ["2025-02", "2025-03", "2025-04"]
    keys_26  = ["2026-02", "2026-03", "2026-04"]

    occ_25  = [by_m[k]["occ"]     for k in keys_25]
    adr_25  = [by_m[k]["adr"]     for k in keys_25]
    rev_25  = [by_m[k]["revpar"]  for k in keys_25]
    rev_25s = [by_m[k]["revenue"] for k in keys_25]

    occ_26  = [by_m[k]["occ"]     for k in keys_26]
    adr_26  = [by_m[k]["adr"]     for k in keys_26]
    rev_26  = [by_m[k]["revpar"]  for k in keys_26]
    rev_26s = [by_m[k]["revenue"] for k in keys_26]

    # Left: summary table
    tbl_data = [
        ["指標",   "2025年", "2026年", "増減"],
        ["OCC（%）",
         f"{sum(occ_25)/3:.1f}%", f"{sum(occ_26)/3:.1f}%",
         f"{sum(occ_26)/3 - sum(occ_25)/3:+.1f}pt"],
        ["ADR（円）",
         f"¥{sum(adr_25)/3:,.0f}", f"¥{sum(adr_26)/3:,.0f}",
         f"{sum(adr_26)/3 - sum(adr_25)/3:+,.0f}円"],
        ["RevPAR（円）",
         f"¥{sum(rev_25)/3:,.0f}", f"¥{sum(rev_26)/3:,.0f}",
         f"{sum(rev_26)/3 - sum(rev_25)/3:+,.0f}円"],
        ["売上（円）",
         f"¥{sum(rev_25s):,}", f"¥{sum(rev_26s):,}",
         f"{sum(rev_26s)-sum(rev_25s):+,}円"],
    ]

    tbl_w = Inches(4.5)
    tbl_h = Inches(5.0)
    tbl = slide.shapes.add_table(
        len(tbl_data), 4,
        Inches(0.3), Inches(0.8), tbl_w, tbl_h).table
    col_widths = [Inches(1.3), Inches(1.05), Inches(1.05), Inches(1.05)]
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = w

    for ri, row in enumerate(tbl_data):
        for ci, val in enumerate(row):
            fg = C_WHITE if ri == 0 else C_DARK
            bg = C_NAVY  if ri == 0 else (C_LBLUE if ri % 2 == 0 else None)
            bold = ri == 0 or ci == 0
            clr = C_RED if (ri > 0 and ci == 3
                            and val.startswith("+")) else fg
            _set_cell(tbl.cell(ri, ci), val, size=10,
                      bold=bold, fg=clr, bg=bg,
                      align=PP_ALIGN.CENTER)

    # Right: 4 grouped bar charts
    x = list(range(3))
    w = 0.32
    fig, axes = plt.subplots(2, 2, figsize=(7.8, 5.5))
    specs = [
        (axes[0, 0], "OCC（稼働率）",      occ_25, occ_26, "%",    True),
        (axes[0, 1], "ADR（平均客室単価）", adr_25, adr_26, "円",   False),
        (axes[1, 0], "RevPAR",             rev_25, rev_26, "円",   False),
        (axes[1, 1], "売上合計",           rev_25s, rev_26s, "円", False),
    ]

    for ax, title, v25, v26, unit, is_pct in specs:
        b1 = ax.bar([xi - w / 2 for xi in x], v25, w,
                    label="2025年", color="#4472C4", zorder=3)
        b2 = ax.bar([xi + w / 2 for xi in x], v26, w,
                    label="2026年", color="#ED7D31", zorder=3)
        ax.set_xticks(x)
        ax.set_xticklabels(months, fontsize=9)
        ax.set_title(title, fontsize=10, fontweight="bold")
        ax.yaxis.grid(True, linestyle="--", alpha=0.4)
        ax.set_axisbelow(True)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.legend(fontsize=8)

        ymax = max(max(v25), max(v26)) * 1.22
        ax.set_ylim(0, 100 if is_pct else ymax)
        for xi, (y1, y2) in enumerate(zip(v25, v26)):
            fmt = (f"{y1:.1f}%", f"{y2:.1f}%") if is_pct else \
                  (f"{y1:,.0f}", f"{y2:,.0f}")
            ax.text(xi - w / 2, y1 * (1.03 if not is_pct else 1) + (1.5 if is_pct else 0),
                    fmt[0], ha="center", fontsize=7.5)
            ax.text(xi + w / 2, y2 * (1.03 if not is_pct else 1) + (1.5 if is_pct else 0),
                    fmt[1], ha="center", fontsize=7.5)

    fig.tight_layout(pad=1.5)
    slide.shapes.add_picture(
        _fig_buf(fig), Inches(4.9), Inches(0.75), Inches(8.2), Inches(6.5))

    _note_box(slide,
              "稼働率はOCC・売上にばらつきがあるものの、ADRは全月で上昇。"
              "単価改善が進んでいます。",
              Inches(0.3), Inches(6.1), Inches(4.4), Inches(0.65))


def slide_operation_change(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_bar(slide, "運用変化")

    before_items = [
        ("01", "毎朝：直近予約・オンハンド\n確認"),
        ("02", "前日比較で料金調整\n（365日先まで）"),
        ("03", "予約内訳を分析"),
        ("04", "セールス・マーケ施策の実施"),
    ]
    after_items = [
        ("01", "毎朝：直近予約・オンハンド\n確認 ＋ TPの動きを確認"),
        ("02", "365日先まで自動調整\n（需要・在庫に応じてリアルタイム調整）"),
        ("03", "予約内訳を多角的に分析\n（サイト別・属性・プラン・部屋）"),
        ("04", "セールス・マーケ施策の実施"),
    ]

    def _box(title, items, left, color):
        box = slide.shapes.add_shape(1, left, Inches(0.85), Inches(5.9), Inches(4.7))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF0, 0xF5, 0xFB)
        box.line.color.rgb = color
        _txb(slide, title, left + Inches(0.15), Inches(0.9),
             Inches(5.6), Inches(0.55),
             size=14, bold=True, color=color)
        for i, (num, text) in enumerate(items):
            y = Inches(1.55) + i * Inches(0.95)
            _txb(slide, num, left + Inches(0.15), y,
                 Inches(0.5), Inches(0.35), size=11, bold=True, color=color)
            sep = slide.shapes.add_shape(1, left + Inches(0.15), y + Inches(0.37),
                                         Inches(5.6), Inches(0.02))
            sep.fill.solid()
            sep.fill.fore_color.rgb = color
            sep.line.fill.background()
            _txb(slide, text, left + Inches(0.15), y + Inches(0.42),
                 Inches(5.6), Inches(0.5),
                 size=11, color=C_DARK, wrap=True)

    _box("【Before】手動オペレーション中心",
         before_items, Inches(0.35), C_NAVY)
    _box("【After】自動化＋判断業務へシフト",
         after_items, Inches(6.95), RGBColor(0xED, 0x7D, 0x31))

    _note_box(slide,
              "価格調整の工数削減により、分析・戦略業務へリソースを集中できる運用へ転換。",
              Inches(0.35), Inches(5.85), Inches(12.6), Inches(0.65))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--revpar1", required=True)
    parser.add_argument("--revpar2", required=True)
    parser.add_argument("--changes", required=True)
    parser.add_argument("--output",  default="sevendays_report.pptx")
    args = parser.parse_args()

    revpar_all = load_revpar(args.revpar1) + load_revpar(args.revpar2)
    seen = {}
    for r in revpar_all:
        seen[r["month"]] = r
    revpar_all = sorted(seen.values(), key=lambda r: r["month"])

    changes_data = load_changes(args.changes)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_intro(prs)
    slide_changes_month(prs, changes_data, 2026, 3)
    slide_changes_month(prs, changes_data, 2026, 4)
    slide_section_break(prs, "昨対比較")
    slide_monthly_comparison(prs, revpar_all)
    slide_operation_change(prs)

    prs.save(args.output)
    print(f"保存しました: {args.output}")


if __name__ == "__main__":
    main()
